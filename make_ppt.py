from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

SLD_LAYOUT_TITLE_AND_CONTENT = 1
SLD_LAYOUT_BLANK = 6
SLD_LAYOUT_PICTURE_WITH_CAPTION = 8

LEFT_OFFSET_INCHES = 0.2
TOP_OFFSET_INCHES = 1
WIDTH_INCHES = 3
HEIGHT_INCHES = 6
HORIZONTAL_PADDING_INCHES = 0.1
VERTICAL_PADDING_INCHES = 0.1
BODY_HEIGHT_INCHES = 5
BODY_WIDTH_INCHES = 3

job_list = [ {'jobID': "0131fe72-09d8-401a-9809-90126e9056a5", 'pictures': ["0131fe72-09d8-401a-9809-90126e9056a5_2be3ad48-42d9-4edb-890b-587d7f2b9cab.jpg", "0131fe72-09d8-401a-9809-90126e9056a5_7daf35c2-e5f7-4fd0-bfa9-f4ca17f856e0.jpg"], 'body': "This is the body"} ]

def create_presentation():
    prs = Presentation()
    return prs

def add_slide(prs, title_string, picture_file, body_string, title_font_size = 18, body_font_size = 18):
    blank_slide_layout = prs.slide_layouts[SLD_LAYOUT_BLANK]
    slide = prs.slides.add_slide(blank_slide_layout)
    # Title
    title = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT_OFFSET_INCHES, Inches(0), Inches(8), Inches(TOP_OFFSET_INCHES))
    title.fill.solid()
    title.fill.fore_color.rgb = RGBColor(255, 255, 255)
    title.line.fill.background()
    title.shadow.inherit = False
    title.text_frame.margin_top = Inches(0)
    title_p = title.text_frame.add_paragraph()
    title_run = title_p.add_run()
    font = title_run.font
    font.name = 'Calibri'
    font.size = Pt(title_font_size)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    font.color.rgb = RGBColor(10, 10, 10)            
    title_run.text = title_string
    # Picture
    im = Image.open(picture_file)
    width, height = im.size
    # let's try to make it full-height
    height_inches = HEIGHT_INCHES
    width_inches = HEIGHT_INCHES * width / height
    # if too large, let's make it fixed-width
    if width_inches > WIDTH_INCHES and body_string:
        height_inches = WIDTH_INCHES * height / width
        width_inches = WIDTH_INCHES
    #pic = slide.shapes.add_picture(pic, Inches(LEFT_OFFSET_INCHES + (WIDTH_INCHES + HORIZONTAL_PADDING_INCHES) * pic_num), Inches(TOP_OFFSET_INCHES), Inches(WIDTH_INCHES), Inches(height_inches))
    slide.shapes.add_picture(picture_file, Inches(LEFT_OFFSET_INCHES), Inches(TOP_OFFSET_INCHES), Inches(width_inches), Inches(height_inches))
    # Body
    if body_string is not None:
        #body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT_OFFSET_INCHES, Inches(TOP_OFFSET_INCHES + max_height_inches), Inches(8), Inches(BODY_HEIGHT_INCHES))
        body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(LEFT_OFFSET_INCHES + width_inches), Inches(TOP_OFFSET_INCHES), Inches(BODY_WIDTH_INCHES), Inches(BODY_HEIGHT_INCHES))
        body.fill.solid()
        body.fill.fore_color.rgb = RGBColor(255, 255, 255)
        body.line.fill.background()
        body.shadow.inherit = False
        body.text_frame.margin_top = Inches(0)
        body_p = body.text_frame.add_paragraph()
        body_p.alignment = PP_ALIGN.LEFT
        body.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        body_run = body_p.add_run()
        font = body_run.font
        font.name = 'Calibri'
        font.size = Pt(body_font_size)
        font.bold = False
        font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        font.color.rgb = RGBColor(10, 10, 10)            
        body_run.text = body_string

def save_presentation(prs, filename):
    prs.save(filename)

def make_ppt(job_list):
    prs = Presentation()
    pic_slide_layout = prs.slide_layouts[SLD_LAYOUT_PICTURE_WITH_CAPTION]
    for job in job_list:
        for pic_num, pic in enumerate(job['pictures']):
            slide = prs.slides.add_slide(pic_slide_layout)
            title = slide.shapes.title
            pic_placeholder = slide.placeholders[1]  # idx key, not position
            pic_placeholder.left = 100
            title.text = "{} {}/{}".format(job['jobID'], pic_num+1, len(job['pictures']))
            picture = pic_placeholder.insert_picture(pic)
            text_placeholder = slide.placeholders[2]
            text_placeholder.text = "this is text"
    prs.save('test.pptx')


def make_ppt2(job_list):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[SLD_LAYOUT_BLANK]
    for job in job_list:
        slide = prs.slides.add_slide(blank_slide_layout)
        title = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, LEFT_OFFSET_INCHES, Inches(0), Inches(8), Inches(TOP_OFFSET_INCHES))
        title.fill.solid()
        title.fill.fore_color.rgb = RGBColor(255, 255, 255)
        title.line.fill.background()
        title.shadow.inherit = False
        title.text_frame.margin_top = Inches(0)
        p = title.text_frame.add_paragraph()
        run = p.add_run()
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(18)
        font.bold = True
        font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        font.color.rgb = RGBColor(10, 10, 10)            
        run.text = job['jobID']
        max_height_inches = 0
        for pic_num, pic in enumerate(job['pictures']):
            #title.text = job['jobID']
            #title.text_frame.font.color.rgb = RGBColor(0x10, 0x10, 0x10)
            im = Image.open(pic)
            width, height = im.size
            height_inches = WIDTH_INCHES*height/width
            if height_inches > max_height_inches:
                max_height_inches = height_inches
            pic = slide.shapes.add_picture(pic, Inches(LEFT_OFFSET_INCHES + (WIDTH_INCHES + HORIZONTAL_PADDING_INCHES) * pic_num), Inches(TOP_OFFSET_INCHES), Inches(WIDTH_INCHES), Inches(height_inches))
        body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, LEFT_OFFSET_INCHES, Inches(TOP_OFFSET_INCHES + max_height_inches), Inches(8), Inches(TOP_OFFSET_INCHES))
        body.fill.solid()
        body.fill.fore_color.rgb = RGBColor(255, 255, 255)
        body.line.fill.background()
        body.shadow.inherit = False
        body.text_frame.margin_top = Inches(0)
        p = body.text_frame.add_paragraph()
        run = p.add_run()
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(18)
        font.bold = False
        font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        font.color.rgb = RGBColor(10, 10, 10)            
        run.text = job['body']
            #title = slide.shapes.title
            #title.text = "{} {}/{}".format(job['jobID'], pic_num+1, len(job['pictures']))
            #text_placeholder = slide.placeholders[2]
            #text_placeholder.text = "this is text"
    prs.save('test.pptx')
